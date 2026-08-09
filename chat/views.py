"""
HTTP surface for standalone chat.

These views do transport only: authenticate, parse, delegate to
`pipeline.run_chat_turn`, serialise. The conversation logic lives in
`pipeline` / `agent`, so the streaming and non-streaming endpoints cannot drift
apart — they are the same call with a different `EventSink`.
"""
from __future__ import annotations

import json
import logging
from typing import Any
from uuid import UUID

from adrf.decorators import api_view
from asgiref.sync import sync_to_async
from django.http import StreamingHttpResponse
from django.shortcuts import get_object_or_404
from django.views.decorators.csrf import csrf_exempt
from rest_framework import status, viewsets
from rest_framework.decorators import (
    api_view as sync_api_view,
    parser_classes,
    permission_classes,
)
from rest_framework.parsers import FormParser, MultiPartParser
from rest_framework.permissions import IsAuthenticated
from rest_framework.response import Response

from workflow_backend.thresholds import (
    DOCUMENT_EXTRACT_CAP,
    IS_LARGE_FILE_THRESHOLD,
    LARGE_FILE_PREVIEW_LENGTH,
)

from .events import Event
from .models import ChatAttachment, ChatMessage, ChatSession
from .pipeline import TurnError, TurnRequest, run_chat_turn
from .serializers import (
    ChatAttachmentSerializer,
    ChatMessageSerializer,
    ChatSessionSerializer,
)
from .sse import SSEBridge, error_frame

logger = logging.getLogger(__name__)


serialize_message = sync_to_async(lambda m: ChatMessageSerializer(m).data)
serialize_attachment = sync_to_async(lambda a: ChatAttachmentSerializer(a).data)


# ── Shared helpers ───────────────────────────────────────────────────────────

async def _get_session(session_id: str, user) -> ChatSession:
    """Fetch the caller's session, or raise `TurnError` with a safe message."""
    try:
        session_uuid = UUID(session_id)
    except (ValueError, AttributeError, TypeError):
        raise TurnError("Invalid session id.") from None

    session = await ChatSession.objects.filter(id=session_uuid, user=user).afirst()
    if session is None:
        raise TurnError("Chat session not found.")
    return session


def _json_body(request) -> dict[str, Any]:
    """Parse a plain-Django request body. Streaming views have no `request.data`."""
    try:
        body = json.loads(request.body or b"{}")
    except (json.JSONDecodeError, UnicodeDecodeError):
        return {}
    return body if isinstance(body, dict) else {}


async def _authenticate(request):
    """
    Resolve the user for a view DRF's decorators cannot wrap.

    `StreamingHttpResponse` bypasses DRF's authentication, so the bearer token
    is verified here instead.
    """
    user = getattr(request, "user", None)
    if user is not None and user.is_authenticated:
        return user

    header = request.META.get("HTTP_AUTHORIZATION", "")
    if not header.startswith("Bearer "):
        return None

    from django.contrib.auth import get_user_model
    from rest_framework_simplejwt.exceptions import TokenError
    from rest_framework_simplejwt.tokens import AccessToken

    try:
        token = AccessToken(header.removeprefix("Bearer ").strip())
        return await get_user_model().objects.aget(id=token["user_id"])
    except (TokenError, KeyError, get_user_model().DoesNotExist) as exc:
        logger.info("[Auth] Rejected streaming request: %s", exc)
        return None


# ── Sessions ─────────────────────────────────────────────────────────────────

class ChatSessionViewSet(viewsets.ModelViewSet):
    """CRUD for standalone chat sessions."""

    serializer_class = ChatSessionSerializer
    permission_classes = [IsAuthenticated]

    def get_queryset(self):
        return ChatSession.objects.filter(user=self.request.user)

    def perform_create(self, serializer) -> None:
        serializer.save(user=self.request.user)

    def perform_destroy(self, instance) -> None:
        """Delete the session along with its RAG documents and vector index."""
        session_id = str(instance.id)

        document_ids = list(
            ChatAttachment.objects
            .filter(session=instance, inference_document__isnull=False)
            .values_list("inference_document_id", flat=True)
        )
        if document_ids:
            try:
                from inference.models import Document

                Document.objects.filter(id__in=document_ids).delete()
            except Exception:
                logger.exception("[Session] RAG cleanup failed for %s", session_id)

        try:
            from inference.engine import get_session_kb_manager

            get_session_kb_manager().clear_session_kb(session_id)
        except Exception:
            logger.exception("[Session] Vector index cleanup failed for %s", session_id)

        instance.delete()


# ── Sending messages ─────────────────────────────────────────────────────────

@api_view(["POST"])
@permission_classes([IsAuthenticated])
async def send_message(request, session_id: str):
    """
    Send a message and wait for the complete reply.

    Same pipeline as the streaming endpoint with events discarded. Prefer
    `send_message_stream` in the UI — this one holds the connection open for the
    whole turn, tool calls included.
    """
    try:
        session = await _get_session(session_id, request.user)
        turn = TurnRequest.parse(request.data)
        outcome = await run_chat_turn(session=session, user=request.user, request=turn)
    except TurnError as exc:
        return Response({"error": str(exc)}, status=status.HTTP_400_BAD_REQUEST)

    return Response({
        "user_message": await serialize_message(outcome.user_message),
        "ai_response": await serialize_message(outcome.assistant_message),
    })


@csrf_exempt
async def send_message_stream(request, session_id: str):
    """
    Send a message and stream the turn as server-sent events.

    Emits `status`, `thinking_chunk`, `content_chunk`, `agent_trace`,
    `sources_update`, `images_update`, `videos_update`, `html_artifact`,
    `ask_permission` and a final `done`; `error` on failure.
    """
    user = await _authenticate(request)
    if user is None:
        return StreamingHttpResponse(
            iter([error_frame("Authentication required.")]),
            content_type="text/event-stream",
        )

    payload = _json_body(request)
    bridge = SSEBridge()

    async def turn() -> None:
        try:
            session = await _get_session(session_id, user)
            outcome = await run_chat_turn(
                session=session,
                user=user,
                request=TurnRequest.parse(payload),
                sink=bridge.sink,
            )
        except TurnError as exc:
            await bridge.emit(Event.ERROR, message=str(exc))
            return

        await bridge.emit(
            Event.DONE,
            user_message=await serialize_message(outcome.user_message),
            ai_response=await serialize_message(outcome.assistant_message),
        )

    response = StreamingHttpResponse(
        bridge.stream(turn()), content_type="text/event-stream"
    )
    response["Cache-Control"] = "no-cache"
    response["X-Accel-Buffering"] = "no"  # let nginx pass chunks straight through
    return response


# ── Deleting messages ────────────────────────────────────────────────────────

def _release_attachment(session: ChatSession, message: ChatMessage) -> None:
    """Delete the attachment a message owns, plus its file and RAG records."""
    raw_id = (message.metadata or {}).get("attachment_id")
    if not raw_id:
        return

    try:
        attachment_id = UUID(raw_id)
    except (ValueError, TypeError):
        return

    attachment = ChatAttachment.objects.filter(id=attachment_id, session=session).first()
    if attachment is None:
        return

    if attachment.file:
        try:
            attachment.file.delete(save=False)
        except OSError as exc:
            logger.warning("[Delete] Could not remove file from disk: %s", exc)

    if attachment.inference_document_id:
        _purge_rag_document(str(session.id), attachment.inference_document_id)

    attachment.delete()


def _purge_rag_document(session_id: str, document_id) -> None:
    """Remove a document from the session vector index and the database."""
    try:
        from asgiref.sync import async_to_sync

        from inference.engine import get_session_knowledge_base
        from inference.models import Document

        try:
            async_to_sync(get_session_knowledge_base(session_id).delete_document)(document_id)
        except Exception:
            # The SQL row must still go even if the index is unavailable —
            # leaving it behind would resurrect the document on reindex.
            logger.warning("[Delete] Vector index removal failed", exc_info=True)

        Document.objects.filter(id=document_id).delete()
    except Exception:
        logger.exception("[Delete] RAG cleanup failed for document %s", document_id)


@sync_api_view(["DELETE"])
@permission_classes([IsAuthenticated])
def delete_message(request, session_id: str, message_id: int):
    """
    Delete a message, or rewind the conversation from it.

    `?rewind=true` removes this message and everything after it;
    `?rewind_after=true` keeps this message and removes what follows.
    """
    try:
        session = ChatSession.objects.filter(id=UUID(session_id), user=request.user).first()
    except (ValueError, TypeError):
        return Response({"error": "Invalid session id."}, status=400)

    if session is None:
        return Response({"error": "Chat session not found."}, status=404)

    message = get_object_or_404(ChatMessage, id=message_id, session=session)

    def flag(name: str) -> bool:
        return request.query_params.get(name, "").lower() == "true"

    if flag("rewind_after"):
        targets = ChatMessage.objects.filter(session=session, id__gt=message_id)
        detail = "Conversation rewound."
    elif flag("rewind"):
        targets = ChatMessage.objects.filter(session=session, id__gte=message_id)
        detail = "Conversation rewound."
    else:
        targets = ChatMessage.objects.filter(id=message.id)
        detail = "Message deleted."

    # Ids are sequential, so `id >=` is both cheaper and more precise than a
    # timestamp comparison when two messages share a creation second.
    for target in targets:
        _release_attachment(session, target)
    targets.delete()

    return Response({"status": detail}, status=status.HTTP_204_NO_CONTENT)


# ── File upload ──────────────────────────────────────────────────────────────

_FILE_TYPES = {
    "image": (".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp"),
    "pdf": (".pdf",),
    "pptx": (".pptx", ".ppt"),
    "text": (".txt", ".md", ".csv", ".json", ".xml", ".html"),
}


def _classify_file(filename: str) -> str:
    lowered = filename.lower()
    for kind, suffixes in _FILE_TYPES.items():
        if lowered.endswith(suffixes):
            return kind
    return "other"


def _extract_pdf_text(data: bytes, max_pages: int = 100) -> str:
    import io

    try:
        from PyPDF2 import PdfReader
    except ImportError:
        return "[PDF extraction needs PyPDF2: pip install PyPDF2]"

    try:
        reader = PdfReader(io.BytesIO(data))
        return "\n\n".join(
            text for page in reader.pages[:max_pages] if (text := page.extract_text())
        )
    except Exception:
        logger.exception("[Upload] PDF extraction failed")
        return ""


def _extract_pptx_text(data: bytes) -> str:
    import io

    try:
        from pptx import Presentation
    except ImportError:
        return "[PPTX extraction needs python-pptx: pip install python-pptx]"

    try:
        return "\n\n".join(
            shape.text
            for slide in Presentation(io.BytesIO(data)).slides
            for shape in slide.shapes
            if getattr(shape, "text", "")
        )
    except Exception:
        logger.exception("[Upload] PPTX extraction failed")
        return ""


def _extract_text(data: bytes, file_type: str) -> str:
    match file_type:
        case "pdf":
            return _extract_pdf_text(data)
        case "pptx":
            return _extract_pptx_text(data)
        case "text":
            return data.decode("utf-8", errors="ignore")
        case _:
            return ""


def _index_for_rag(request, upload, attachment: ChatAttachment, text: str) -> None:
    """Register the file in the user knowledge base and index it in background."""
    import threading

    from inference.models import Document
    from inference.tasks import process_document

    document = Document.objects.create(
        user=request.user, name=upload.name, content_text=text, file=upload,
        file_type=attachment.file_type, file_size=upload.size, status="pending",
    )
    attachment.inference_document = document
    attachment.save(update_fields=["inference_document"])

    # Indexed in-process: this deployment has no Celery broker.
    threading.Thread(target=process_document, args=(document.id,), daemon=True).start()


@sync_api_view(["POST"])
@permission_classes([IsAuthenticated])
@parser_classes([MultiPartParser, FormParser])
def upload_file(request, session_id: str):
    """Attach a file to a session, extracting its text for context and RAG."""
    try:
        session = ChatSession.objects.filter(id=UUID(session_id), user=request.user).first()
    except (ValueError, TypeError):
        return Response({"error": "Invalid session id."}, status=400)

    if session is None:
        return Response({"error": "Chat session not found."}, status=404)

    upload = request.FILES.get("file")
    if not upload:
        return Response({"error": "No file provided."}, status=400)

    file_type = _classify_file(upload.name)
    data = upload.read()
    upload.seek(0)
    text = _extract_text(data, file_type)

    attachment = ChatAttachment.objects.create(
        session=session,
        file=upload,
        filename=upload.name,
        file_type=file_type,
        file_size=upload.size,
        extracted_text=text[:DOCUMENT_EXTRACT_CAP],
        is_large_file=len(text) > IS_LARGE_FILE_THRESHOLD,
    )

    if file_type in ("pdf", "pptx", "text") and text:
        try:
            _index_for_rag(request, upload, attachment, text)
        except Exception:
            logger.exception("[Upload] RAG indexing failed for %s", upload.name)

    preview = (
        f"\n\nContext summary ({len(text)} chars):\n{text[:LARGE_FILE_PREVIEW_LENGTH]}..."
        if text else ""
    )
    ChatMessage.objects.create(
        session=session,
        role="system",
        content=f"📎 File uploaded: **{upload.name}** ({file_type}, {upload.size} bytes){preview}",
        message_type="system",
        metadata={
            "attachment_id": str(attachment.id),
            "file_type": file_type,
            "has_extracted_text": bool(text),
        },
    )

    return Response({
        "attachment": ChatAttachmentSerializer(attachment).data,
        "extracted_text_length": len(text),
        "message": f'File "{upload.name}" uploaded successfully.',
    })


# ── Direct tool execution ────────────────────────────────────────────────────

@api_view(["POST"])
@permission_classes([IsAuthenticated])
async def execute_tool_view(request):
    """
    Run a single tool directly, for UI affordances that bypass the agent.

    Sensitive tools are refused here: they are gated by human approval inside
    the agent loop, and reaching them through this endpoint would be a way
    around that gate rather than a shortcut to it.
    """
    from .tools import SENSITIVE_TOOLS, execute_tool

    name = request.data.get("tool")
    args = request.data.get("args")

    if not name:
        return Response({"error": "'tool' is required."}, status=400)
    if not isinstance(args, dict):
        return Response({"error": "'args' must be an object."}, status=400)
    if name in SENSITIVE_TOOLS:
        return Response(
            {
                "error": "This tool requires human approval and cannot be run directly.",
                "status": "blocked",
            },
            status=403,
        )

    try:
        raw = await execute_tool(name, args, {"user_id": request.user.id})
    except Exception as exc:
        logger.exception("[Tools] Direct execution of %s failed", name)
        return Response({"error": str(exc)}, status=500)

    try:
        return Response(json.loads(raw))
    except (json.JSONDecodeError, TypeError):
        return Response({"result": raw})
