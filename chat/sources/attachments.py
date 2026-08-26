"""
File and RAG lifecycle for chat attachments.

Upload classifies a file, pulls its text out, and registers it in the user
knowledge base; deletion has to undo all three — the file on disk, the vector
index entry, and the `Document` row. Both halves live here so they cannot drift:
anything the upload path creates, the release path knows how to remove.
"""
from __future__ import annotations

import logging
from uuid import UUID

from chat.models import ChatAttachment, ChatMessage, ChatSession

logger = logging.getLogger(__name__)


# ── Classification and text extraction ───────────────────────────────────────

FILE_TYPES = {
    "image": (".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp"),
    "pdf": (".pdf",),
    "pptx": (".pptx", ".ppt"),
    "text": (".txt", ".md", ".csv", ".json", ".xml", ".html"),
}


def classify_file(filename: str) -> str:
    lowered = filename.lower()
    for kind, suffixes in FILE_TYPES.items():
        if lowered.endswith(suffixes):
            return kind
    return "other"


def extract_pdf_text(data: bytes, max_pages: int = 100) -> str:
    import io

    try:
        from PyPDF2 import PdfReader
    except ImportError:
        return "[PDF extraction needs PyPDF2: pip install PyPDF2]"

    try:
        reader = PdfReader(io.BytesIO(data))
        return "\n\n".join(
            text for page in reader.pages[:max_pages] if (text := page.extract_text())
        )
    except Exception:
        logger.exception("[Upload] PDF extraction failed")
        return ""


def extract_pptx_text(data: bytes) -> str:
    import io

    try:
        from pptx import Presentation
    except ImportError:
        return "[PPTX extraction needs python-pptx: pip install python-pptx]"

    try:
        return "\n\n".join(
            shape.text
            for slide in Presentation(io.BytesIO(data)).slides
            for shape in slide.shapes
            if getattr(shape, "text", "")
        )
    except Exception:
        logger.exception("[Upload] PPTX extraction failed")
        return ""


def extract_text(data: bytes, file_type: str) -> str:
    match file_type:
        case "pdf":
            return extract_pdf_text(data)
        case "pptx":
            return extract_pptx_text(data)
        case "text":
            return data.decode("utf-8", errors="ignore")
        case _:
            return ""


# ── Indexing ─────────────────────────────────────────────────────────────────

def index_for_rag(user, upload, attachment: ChatAttachment, text: str) -> None:
    """Register the file in the user knowledge base and index it in background."""
    import threading

    from inference.models import Document
    from inference.tasks import process_document
    from inference.utils import normalize_file_type

    # `ChatAttachment.file_type` speaks its own five-value vocabulary
    # ('image'/'pdf'/'pptx'/'text'/'other'), which is not Document's — 'pptx',
    # 'text' and 'other' are not FILE_TYPE_CHOICES at all. Normalising from the
    # filename keeps every producer of a Document row agreeing with every
    # consumer of `Document.file_type` (text extraction, vision extraction).
    document = Document.objects.create(
        user=user, name=upload.name, content_text=text, file=upload,
        file_type=normalize_file_type(upload.name, getattr(upload, 'content_type', '')),
        file_size=upload.size, status="pending",
    )
    attachment.inference_document = document
    attachment.save(update_fields=["inference_document"])

    # Indexed in-process: this deployment has no Celery broker.
    threading.Thread(target=process_document, args=(document.id,), daemon=True).start()


# ── Release ──────────────────────────────────────────────────────────────────

def release_attachment(session: ChatSession, message: ChatMessage) -> None:
    """Delete the attachment a message owns, plus its file and RAG records."""
    raw_id = (message.metadata or {}).get("attachment_id")
    if not raw_id:
        return

    try:
        attachment_id = UUID(raw_id)
    except (ValueError, TypeError):
        return

    attachment = ChatAttachment.objects.filter(id=attachment_id, session=session).first()
    if attachment is None:
        return

    if attachment.file:
        try:
            attachment.file.delete(save=False)
        except OSError as exc:
            logger.warning("[Delete] Could not remove file from disk: %s", exc)

    if attachment.inference_document_id:
        purge_rag_document(str(session.id), attachment.inference_document_id)

    attachment.delete()


def purge_rag_document(session_id: str, document_id) -> None:
    """Remove a document from the session vector index and the database."""
    try:
        from asgiref.sync import async_to_sync

        from inference.engine import get_session_knowledge_base
        from inference.models import Document

        try:
            async_to_sync(get_session_knowledge_base(session_id).delete_document)(document_id)
        except Exception:
            # The SQL row must still go even if the index is unavailable —
            # leaving it behind would resurrect the document on reindex.
            logger.warning("[Delete] Vector index removal failed", exc_info=True)

        Document.objects.filter(id=document_id).delete()
    except Exception:
        logger.exception("[Delete] RAG cleanup failed for document %s", document_id)


def purge_session(session: ChatSession) -> None:
    """
    Drop everything a session's attachments left outside the session's own rows.

    Cascading the delete would take the `ChatAttachment` rows and stop there,
    leaving the `Document` rows and the vector index behind — orphaned entries
    that keep answering retrieval queries for a conversation that no longer
    exists.
    """
    session_id = str(session.id)

    document_ids = list(
        ChatAttachment.objects
        .filter(session=session, inference_document__isnull=False)
        .values_list("inference_document_id", flat=True)
    )
    if document_ids:
        try:
            from inference.models import Document

            Document.objects.filter(id__in=document_ids).delete()
        except Exception:
            logger.exception("[Session] RAG cleanup failed for %s", session_id)

    try:
        from inference.engine import clear_session_kb

        clear_session_kb(session_id)
    except Exception:
        logger.exception("[Session] Vector index cleanup failed for %s", session_id)
